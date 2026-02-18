"""
PowerPoint Presentation Generator for Cloud Removal AI v5.0
Creates a professional .pptx presentation with all project features
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    """Generate the complete PowerPoint presentation"""
    
    # Create presentation object
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Define color scheme
    PRIMARY_COLOR = RGBColor(102, 126, 234)  # #667eea
    SECONDARY_COLOR = RGBColor(118, 75, 162)  # #764ba2
    WHITE = RGBColor(255, 255, 255)
    DARK_GRAY = RGBColor(51, 51, 51)
    
    # Slide 1: Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add gradient background (simulated with rectangle)
    left = top = Inches(0)
    width = prs.slide_width
    height = prs.slide_height
    shape = slide.shapes.add_shape(1, left, top, width, height)
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = PRIMARY_COLOR
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "🛰️ Cloud Removal AI v5.0"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(60)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE
    title_para.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(8), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Advanced Satellite Imagery Processing"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(32)
    subtitle_para.font.color.rgb = WHITE
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    # Version
    version_box = slide.shapes.add_textbox(Inches(3), Inches(5), Inches(4), Inches(0.6))
    version_frame = version_box.text_frame
    version_frame.text = "High-Resolution Cloud Removal with AI"
    version_para = version_frame.paragraphs[0]
    version_para.font.size = Pt(20)
    version_para.font.color.rgb = WHITE
    version_para.alignment = PP_ALIGN.CENTER
    
    # Slide 2: Project Overview
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title only
    title = slide.shapes.title
    title.text = "📋 Project Overview"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    
    # Content
    left = Inches(1)
    top = Inches(2)
    width = Inches(8)
    height = Inches(1.5)
    
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "An advanced AI-powered application for removing clouds from satellite imagery using Deep Learning and Computer Vision techniques."
    p.font.size = Pt(22)
    p.alignment = PP_ALIGN.CENTER
    
    # Stats boxes
    stats = [
        ("10+", "Features"),
        ("5", "LULC Classes"),
        ("3D", "Terrain Viz")
    ]
    
    box_width = Inches(2.3)
    box_height = Inches(1.5)
    start_left = Inches(1.2)
    top = Inches(4.2)
    
    for i, (number, label) in enumerate(stats):
        left = start_left + i * Inches(2.8)
        shape = slide.shapes.add_shape(1, left, top, box_width, box_height)
        fill = shape.fill
        fill.solid()
        fill.fore_color.rgb = PRIMARY_COLOR
        
        # Number
        text_frame = shape.text_frame
        text_frame.clear()
        p = text_frame.paragraphs[0]
        p.text = number
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        
        # Label
        p = text_frame.add_paragraph()
        p.text = label
        p.font.size = Pt(18)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
    
    # Slide 3: Core Features
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "⚡ Core Features"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    
    features = [
        ("☁️ Cloud Removal", "GAN-based UNet architecture"),
        ("🏔️ 3D Terrain", "Interactive visualization"),
        ("🏘️ LULC Classification", "5-class land cover"),
        ("🤖 AI Chatbot", "Tamil/Tanglish support"),
        ("🌱 Vegetation Analysis", "VARI health index"),
        ("🛣️ Infrastructure", "Road extraction"),
        ("📊 Analytics", "Comprehensive dashboard"),
        ("📦 Batch Processing", "Multiple images")
    ]
    
    box_width = Inches(3.5)
    box_height = Inches(1.2)
    
    for i, (title_text, desc) in enumerate(features):
        row = i // 2
        col = i % 2
        left = Inches(1 + col * 4)
        top = Inches(2 + row * 1.4)
        
        shape = slide.shapes.add_shape(1, left, top, box_width, box_height)
        fill = shape.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(245, 247, 250)
        
        text_frame = shape.text_frame
        text_frame.clear()
        text_frame.margin_left = Inches(0.2)
        text_frame.margin_top = Inches(0.15)
        
        p = text_frame.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = PRIMARY_COLOR
        
        p = text_frame.add_paragraph()
        p.text = desc
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(85, 85, 85)
    
    # Slide 4: What's New in v5.0
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "🎉 What's New in v5.0"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    
    # Content
    left = Inches(1.5)
    top = Inches(2)
    width = Inches(7)
    height = Inches(4)
    
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    
    # Section 1
    p = tf.paragraphs[0]
    p.text = "🏔️ 3D Terrain Reconstruction"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = SECONDARY_COLOR
    p.space_after = Pt(10)
    
    items = [
        "Interactive 3D visualization with Plotly",
        "Adjustable height exaggeration (0.5x - 5x)",
        "Two estimation methods: Brightness & Gradient",
        "Full camera controls: Rotate, zoom, pan"
    ]
    
    for item in items:
        p = tf.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(18)
        p.level = 0
        p.space_after = Pt(8)
    
    # Section 2
    p = tf.add_paragraph()
    p.text = "\n🏘️ Land Use & Land Cover Classification"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = SECONDARY_COLOR
    p.space_before = Pt(15)
    p.space_after = Pt(10)
    
    items2 = [
        "5 Classes: Water, Forest, Urban, Barren, Vegetation",
        "Color-coded maps with interactive legend",
        "Pie charts showing distribution percentages"
    ]
    
    for item in items2:
        p = tf.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(18)
        p.level = 0
        p.space_after = Pt(8)
    
    # Slide 5: AI Chatbot
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "🤖 AI Chatbot Assistant"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    
    # Left column
    left_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(4), Inches(4))
    tf = left_box.text_frame
    
    p = tf.paragraphs[0]
    p.text = "Features"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = SECONDARY_COLOR
    p.space_after = Pt(12)
    
    features_list = [
        "Natural language queries",
        "Tamil/Tanglish support",
        "Chat history tracking",
        "Real-time map analysis"
    ]
    
    for item in features_list:
        p = tf.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(20)
        p.space_after = Pt(10)
    
    # Right column
    right_box = slide.shapes.add_textbox(Inches(5.5), Inches(2), Inches(4), Inches(4))
    tf = right_box.text_frame
    
    p = tf.paragraphs[0]
    p.text = "Example Queries"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = SECONDARY_COLOR
    p.space_after = Pt(12)
    
    queries = [
        '"How much water is in this area?"',
        '"Where is the forest located?"',
        '"Is there more urban or vegetation?"',
        '"Evlo thanni irukku?" (Tamil)',
        '"Tell me about this map"'
    ]
    
    for query in queries:
        p = tf.add_paragraph()
        p.text = "• " + query
        p.font.size = Pt(16)
        p.space_after = Pt(8)
    
    # Slide 6: Technical Architecture
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "🏗️ Technical Architecture"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    
    # Flow diagram
    flow_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(2.5))
    tf = flow_box.text_frame
    
    flow_steps = [
        "Input → Cloudy Satellite Image + SAR Data",
        "Processing → GAN-based UNet (Generator + Discriminator)",
        "Cloud Removal → Clean Image Output",
        "Analysis → 3D Terrain + LULC + Vegetation + Infrastructure",
        "Output → Interactive Dashboard + PDF Report"
    ]
    
    for step in flow_steps:
        p = tf.add_paragraph() if tf.paragraphs[0].text else tf.paragraphs[0]
        p.text = step
        p.font.size = Pt(18)
        p.space_after = Pt(12)
    
    # Model details
    details_box = slide.shapes.add_textbox(Inches(1), Inches(4.8), Inches(8), Inches(2))
    tf = details_box.text_frame
    
    p = tf.paragraphs[0]
    p.text = "Model Details"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = SECONDARY_COLOR
    p.space_after = Pt(10)
    
    details = [
        "Architecture: UNet with skip connections",
        "Training: GAN (Generative Adversarial Network)",
        "Input Channels: 4 (RGB + SAR)",
        "Output: 3 (RGB clean image)"
    ]
    
    for detail in details:
        p = tf.add_paragraph()
        p.text = "• " + detail
        p.font.size = Pt(18)
        p.space_after = Pt(6)
    
    # Slide 7: Technology Stack
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "🛠️ Technology Stack"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    
    tech_categories = [
        ("Deep Learning & AI", ["PyTorch", "torchvision", "Google Gemini AI"]),
        ("Computer Vision", ["OpenCV", "PIL/Pillow", "scikit-image"]),
        ("Visualization & UI", ["Streamlit", "Plotly", "Matplotlib"]),
        ("Data Processing", ["NumPy", "Pandas", "SciPy"])
    ]
    
    top = Inches(2)
    for category, techs in tech_categories:
        # Category title
        cat_box = slide.shapes.add_textbox(Inches(1), top, Inches(8), Inches(0.5))
        tf = cat_box.text_frame
        p = tf.paragraphs[0]
        p.text = category
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = SECONDARY_COLOR
        
        # Tech badges
        badge_top = top + Inches(0.6)
        for i, tech in enumerate(techs):
            left = Inches(1.5 + i * 2.3)
            shape = slide.shapes.add_shape(1, left, badge_top, Inches(2), Inches(0.5))
            fill = shape.fill
            fill.solid()
            fill.fore_color.rgb = PRIMARY_COLOR
            
            text_frame = shape.text_frame
            text_frame.clear()
            p = text_frame.paragraphs[0]
            p.text = tech
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER
        
        top += Inches(1.3)
    
    # Slide 8: How to Run
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "🚀 How to Run"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    
    # Local installation
    local_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(2))
    tf = local_box.text_frame
    
    p = tf.paragraphs[0]
    p.text = "Local Installation"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = SECONDARY_COLOR
    p.space_after = Pt(12)
    
    commands = [
        'cd "path/to/project"',
        'pip install -r requirements.txt',
        'streamlit run src/app.py'
    ]
    
    for cmd in commands:
        p = tf.add_paragraph()
        p.text = cmd
        p.font.size = Pt(16)
        p.font.name = 'Courier New'
        p.space_after = Pt(8)
    
    # Google Colab
    colab_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(2))
    tf = colab_box.text_frame
    
    p = tf.paragraphs[0]
    p.text = "Google Colab"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = SECONDARY_COLOR
    p.space_after = Pt(12)
    
    steps = [
        "1. Upload project to Google Drive",
        "2. Open Cloud_Removal_Colab.ipynb",
        "3. Mount Drive and install dependencies",
        "4. Configure ngrok for public URL",
        "5. Run the app!"
    ]
    
    for step in steps:
        p = tf.add_paragraph()
        p.text = step
        p.font.size = Pt(18)
        p.space_after = Pt(6)
    
    # Slide 9: Use Cases
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "🎯 Use Cases & Applications"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    
    use_cases = [
        ("🌾 Agriculture", "Crop monitoring and yield prediction"),
        ("🏙️ Urban Planning", "Infrastructure mapping"),
        ("🌊 Environmental", "Water detection and monitoring"),
        ("🚨 Disaster Response", "Flood mapping and assessment"),
        ("🗺️ Cartography", "Map creation and terrain analysis"),
        ("🔬 Research", "Academic and geospatial studies")
    ]
    
    box_width = Inches(3.5)
    box_height = Inches(1.3)
    
    for i, (title_text, desc) in enumerate(use_cases):
        row = i // 2
        col = i % 2
        left = Inches(1 + col * 4)
        top = Inches(2 + row * 1.5)
        
        shape = slide.shapes.add_shape(1, left, top, box_width, box_height)
        fill = shape.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(245, 247, 250)
        
        text_frame = shape.text_frame
        text_frame.clear()
        text_frame.margin_left = Inches(0.2)
        text_frame.margin_top = Inches(0.2)
        
        p = text_frame.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = PRIMARY_COLOR
        
        p = text_frame.add_paragraph()
        p.text = desc
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(85, 85, 85)
    
    # Slide 10: Thank You
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Background
    shape = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = PRIMARY_COLOR
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Thank You! 🙏"
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(8), Inches(0.6))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Cloud Removal AI v5.0"
    p.font.size = Pt(32)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Footer
    footer_box = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(1))
    tf = footer_box.text_frame
    p = tf.paragraphs[0]
    p.text = "🛰️ Advanced Satellite Imagery Processing"
    p.font.size = Pt(24)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    p = tf.add_paragraph()
    p.text = "\nQuestions?"
    p.font.size = Pt(28)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Save presentation
    output_path = "Cloud_Removal_AI_v5_Presentation.pptx"
    prs.save(output_path)
    print(f"✅ Presentation created successfully: {output_path}")
    print(f"📊 Total slides: {len(prs.slides)}")
    return output_path

if __name__ == "__main__":
    import sys
    import io
    
    # Fix encoding for Windows console
    if sys.platform == 'win32':
        sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')
    
    print("🎨 Generating PowerPoint presentation...")
    print("=" * 60)
    output_file = create_presentation()
    print("=" * 60)
    print(f"✨ Done! Open '{output_file}' to view your presentation.")
